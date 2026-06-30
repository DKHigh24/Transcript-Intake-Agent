"""
presentation_builder.py
Generates the upcoming session PPTX meeting deck from accumulated opportunity history.

Called automatically from run_weekly() after the trend pipeline completes.
Output: output/meeting_presentations/Session_<N>_<M>_<D>_<YYYY>.pptx

Slide structure (mirrors Sessions 2–5):
  1.  Title slide
  2.  Confidential / NYSE placeholder
  3.  Agenda
  4.  What we aligned on (prior session recap, LLM-assisted)
  5.  Demo backlog (top opportunities as candidates)
  6.  Demo placeholder (update before meeting)
  7+. KB Article Candidate(s) — one slide per active candidate (2+ sessions)
  Last. Demo Intake / KB Strategy / Next Steps
"""

from __future__ import annotations

import json
import traceback
from datetime import date, timedelta
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Color palette (Acuity brand) ──────────────────────────────────────────────
NAVY       = RGBColor(0x1F, 0x38, 0x64)
LIGHT_NAVY = RGBColor(0x3A, 0x4B, 0x5C)
TEAL       = RGBColor(0x47, 0xA8, 0xBD)
GOLD       = RGBColor(0xE8, 0xA8, 0x38)
GREEN      = RGBColor(0x5D, 0xB7, 0xA0)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x2C, 0x3E, 0x50)
MUTED      = RGBColor(0x6B, 0x7A, 0x8D)
LIGHT      = RGBColor(0xF4, 0xF7, 0xFB)
BORDER     = RGBColor(0xDD, 0xE3, 0xED)
ROW_ALT    = RGBColor(0xF7, 0xF9, 0xFC)

# ── Layout constants ───────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

PRES_DIR = Path("output") / "meeting_presentations"

COMPANY = (
    "Eureka\u00ae, nLight\u00ae, Atrius\u00ae, Distech Controls:  "
    "Acuity Inc. \u2014 Corporate Headquarters"
)
GROUP = "Electronics AI Working Group"


# ── Low-level shape helpers ───────────────────────────────────────────────────

def _blank_slide(prs: Presentation):
    """Return a new blank slide using the layout with fewest placeholders."""
    best = prs.slide_layouts[6]
    min_ph = len(list(best.placeholders))
    for layout in prs.slide_layouts:
        ph_count = len(list(layout.placeholders))
        if ph_count < min_ph:
            min_ph = ph_count
            best = layout
    return prs.slides.add_slide(best)


def _slide_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def _rect(slide, x, y, w, h, fill: RGBColor | None = None,
          border: RGBColor | None = None):
    """Add a rectangle shape, optionally filled and/or bordered."""
    shape = slide.shapes.add_shape(1, x, y, w, h)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border is not None:
        shape.line.color.rgb = border
    else:
        shape.line.fill.background()
    return shape


def _text_frame(
    shape,
    lines: list[tuple[str, int, bool, RGBColor, PP_ALIGN]],
    anchor=MSO_ANCHOR.TOP,
) -> None:
    """
    Populate a shape's text frame.
    Each element in `lines` is (text, font_pt, bold, color, alignment).
    """
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (text, pt, bold, color, align) in enumerate(lines):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(pt)
        run.font.bold = bold
        run.font.color.rgb = color


def _bullet_frame(
    shape,
    bullets: list[str],
    pt: int,
    color: RGBColor,
    prefix: str = "\u2022  ",
    gap_pt: int = 4,
) -> None:
    """Populate a shape's text frame with bullet lines, with a small gap between each."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            # gap paragraph
            gap_p = tf.add_paragraph()
            gap_run = gap_p.add_run()
            gap_run.text = ""
            gap_run.font.size = Pt(gap_pt)
            gap_run.font.color.rgb = color
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{prefix}{bullet}"
        run.font.size = Pt(pt)
        run.font.color.rgb = color


# ── Common header bar ─────────────────────────────────────────────────────────

def _header_bar(slide, title: str, subtitle: str = "") -> None:
    """Full-width dark navy header bar with white title and optional subtitle."""
    bar_h = Inches(1.4) if subtitle else Inches(1.2)
    _rect(slide, Inches(0), Inches(0), W, bar_h, fill=NAVY)
    sh = _rect(slide, Inches(0.45), Inches(0), W - Inches(0.9), bar_h)
    lines = [(title, 24, True, WHITE, PP_ALIGN.LEFT)]
    if subtitle:
        lines.append((subtitle, 12, False, RGBColor(0xB9, 0xC4, 0xD1), PP_ALIGN.LEFT))
    _text_frame(sh, lines, MSO_ANCHOR.MIDDLE)


# ── Individual slide builders ─────────────────────────────────────────────────

def _slide_title(prs: Presentation, session_num: int, next_date: date) -> None:
    """Slide 1 – Title slide (dark navy background)."""
    slide = _blank_slide(prs)
    _slide_bg(slide, NAVY)

    # Teal accent strip (left edge)
    _rect(slide, Inches(0), Inches(0), Inches(0.22), H, fill=TEAL)

    # Top branding stripe
    _rect(slide, Inches(0.22), Inches(0), W - Inches(0.22), Inches(1.05), fill=LIGHT_NAVY)

    # Company name
    sh_co = _rect(slide, Inches(0.5), Inches(0.1), W - Inches(1.0), Inches(0.85))
    _text_frame(sh_co, [(COMPANY, 11, False, RGBColor(0xB9, 0xC4, 0xD1), PP_ALIGN.LEFT)],
                MSO_ANCHOR.MIDDLE)

    # Working group label
    sh_grp = _rect(slide, Inches(0.5), Inches(1.3), W - Inches(1.0), Inches(0.5))
    _text_frame(sh_grp, [(GROUP, 15, False, TEAL, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)

    # Session title
    sh_title = _rect(slide, Inches(0.5), Inches(1.9), W - Inches(1.0), Inches(2.0))
    _text_frame(sh_title,
                [(f"Session {session_num}  \u2014  Upcoming Discussion & Demo", 34, True, WHITE, PP_ALIGN.LEFT)],
                MSO_ANCHOR.MIDDLE)

    # Date
    sh_date = _rect(slide, Inches(0.5), Inches(4.1), Inches(6.5), Inches(0.65))
    _text_frame(sh_date,
                [(next_date.strftime("%B %d, %Y"), 22, False, GOLD, PP_ALIGN.LEFT)],
                MSO_ANCHOR.MIDDLE)

    # Footer / NYSE
    sh_foot = _rect(slide, Inches(0.5), H - Inches(0.75), W - Inches(1.0), Inches(0.55))
    _text_frame(sh_foot,
                [("Proprietary & Confidential  |  NYSE: AYI  |  Draft — for internal review only",
                  10, False, MUTED, PP_ALIGN.LEFT)],
                MSO_ANCHOR.MIDDLE)


def _slide_confidential(prs: Presentation, slide_num: int) -> None:
    """Slide 2 – Proprietary & Confidential placeholder."""
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)

    sh = _rect(slide, Inches(1.5), Inches(2.3), Inches(10.33), Inches(3.0),
               fill=WHITE, border=BORDER)
    _text_frame(sh, [
        ("Proprietary & Confidential", 30, True, NAVY, PP_ALIGN.CENTER),
        ("NYSE: AYI", 18, False, MUTED, PP_ALIGN.CENTER),
        ("", 10, False, WHITE, PP_ALIGN.CENTER),
        (str(slide_num), 14, False, MUTED, PP_ALIGN.CENTER),
    ], MSO_ANCHOR.MIDDLE)


def _slide_agenda(prs: Presentation, agenda_items: list[str], session_num: int) -> None:
    """Slide 3 – Agenda."""
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, "AGENDA", f"Session {session_num}")

    sh_label = _rect(slide, Inches(0.5), Inches(1.55), Inches(5.0), Inches(0.45))
    _text_frame(sh_label, [("Today\u2019s Discussion", 14, True, NAVY, PP_ALIGN.LEFT)])

    sh_bullets = _rect(slide, Inches(0.5), Inches(2.1), Inches(12.0), Inches(5.0),
                       fill=WHITE, border=BORDER)
    _bullet_frame(sh_bullets, agenda_items, 14, DARK)


def _slide_aligned_on(prs: Presentation, bullets: list[str], prev_date: str) -> None:
    """Slide 4 – What we aligned on (prior session recap)."""
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, "What We Aligned On", f"Recap from Session — {prev_date}")

    # Left green accent
    _rect(slide, Inches(0), Inches(1.4), Inches(0.18), H - Inches(1.4), fill=GREEN)

    sh = _rect(slide, Inches(0.45), Inches(1.6), Inches(12.45), Inches(5.55),
               fill=WHITE, border=BORDER)
    _bullet_frame(sh, bullets, 14, DARK)


def _slide_demo_backlog(prs: Presentation, demos: list[dict]) -> None:
    """Slide 5 – Demo Backlog table."""
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, "Demo Backlog")

    col_headers = ["Presenter / SME", "Opportunity / Topic", "Signal", "Status"]
    col_x      = [Inches(0.35), Inches(2.85), Inches(9.35), Inches(11.2)]
    col_widths  = [Inches(2.4),  Inches(6.4),  Inches(1.75), Inches(1.8)]
    row_h = Inches(0.43)
    header_y = Inches(1.55)

    # Header row
    _rect(slide, Inches(0.35), header_y, W - Inches(0.7), row_h, fill=NAVY)
    for hdr, cx, cw in zip(col_headers, col_x, col_widths):
        sh = _rect(slide, cx, header_y, cw, row_h)
        _text_frame(sh, [(hdr, 11, True, WHITE, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)

    # Data rows
    for r_idx, demo in enumerate(demos[:8]):
        y = header_y + row_h * (r_idx + 1)
        bg = WHITE if r_idx % 2 == 0 else ROW_ALT
        _rect(slide, Inches(0.35), y, W - Inches(0.7), row_h, fill=bg)
        row_vals = [
            (demo.get("owner") or "TBD")[:30],
            (demo.get("title") or "")[:65],
            (demo.get("signal") or "")[:22],
            demo.get("status", "Queued"),
        ]
        for val, cx, cw in zip(row_vals, col_x, col_widths):
            sh = _rect(slide, cx, y, cw, row_h)
            _text_frame(sh, [(val, 11, False, DARK, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)


def _slide_demo(prs: Presentation, session_num: int) -> None:
    """Slide 6 – Demo placeholder (update before the meeting)."""
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, f"Demo \u2013 Session {session_num}")
    _rect(slide, Inches(0), Inches(1.4), Inches(0.18), H - Inches(1.4), fill=TEAL)

    sh = _rect(slide, Inches(1.5), Inches(2.1), Inches(10.33), Inches(4.0),
               fill=WHITE, border=BORDER)
    _text_frame(sh, [
        ("Presenter TBD \u2014 Confirm Before Meeting", 24, True, MUTED, PP_ALIGN.CENTER),
        ("", 10, False, WHITE, PP_ALIGN.CENTER),
        ("Update this slide with presenter name and demo topic.", 14, False, DARK, PP_ALIGN.CENTER),
        ("Rough work is welcome \u2014 demos do not need to be polished.", 13, False, MUTED, PP_ALIGN.CENTER),
    ], MSO_ANCHOR.MIDDLE)


def _slide_kb_candidate(prs: Presentation, idx: int, entry: dict) -> None:
    """KB Article Candidate slide — one per candidate."""
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)

    last_occ = entry["occurrences"][-1] if entry.get("occurrences") else {}
    owner = (
        last_occ.get("SuggestedSMEChampionText")
        or last_occ.get("SuggestedBusinessOwnerText")
        or "TBD"
    )
    tool     = last_occ.get("PrimaryTool", "") or "TBD"
    bucket   = last_occ.get("OperatingBucket", "") or entry.get("bucket", "") or "TBD"
    signal   = last_occ.get("SignalStrength", "") or "TBD"
    next_step = (last_occ.get("NextStep") or "Identify owner and draft KB article.").strip()
    evidence  = (last_occ.get("EvidenceSummary") or
                 "Recurring pattern across multiple sessions.").strip()
    ai_type   = last_occ.get("AIUseCaseType", "") or "TBD"
    occ_count = len(entry.get("occurrences", []))

    _header_bar(slide,
                f"KB Article Candidate #{idx}  |  {entry['title'][:55]}",
                "Why it should become a KB article")

    left_x, right_x = Inches(0.4), Inches(6.85)
    col_w = Inches(6.2)
    content_y = Inches(1.55)

    # Left: metadata card
    sh_left = _rect(slide, left_x, content_y, col_w, Inches(5.7), fill=WHITE, border=BORDER)
    _text_frame(sh_left, [
        (f"Owner / SME:         {owner}", 13, False, DARK, PP_ALIGN.LEFT),
        ("", 5, False, WHITE, PP_ALIGN.LEFT),
        (f"Type:                {ai_type}", 13, False, DARK, PP_ALIGN.LEFT),
        ("", 5, False, WHITE, PP_ALIGN.LEFT),
        (f"Primary Tools:       {tool}", 13, False, DARK, PP_ALIGN.LEFT),
        ("", 5, False, WHITE, PP_ALIGN.LEFT),
        (f"Operating Bucket:    {bucket}", 13, False, DARK, PP_ALIGN.LEFT),
        ("", 5, False, WHITE, PP_ALIGN.LEFT),
        (f"Signal:              {signal}", 13, False, DARK, PP_ALIGN.LEFT),
        ("", 5, False, WHITE, PP_ALIGN.LEFT),
        (f"Sessions Seen:       {occ_count}", 13, True, NAVY, PP_ALIGN.LEFT),
        (f"First Seen:          {entry.get('first_seen', 'N/A')}", 13, False, NAVY, PP_ALIGN.LEFT),
        (f"Last Seen:           {entry.get('last_seen', 'N/A')}", 13, False, NAVY, PP_ALIGN.LEFT),
    ], MSO_ANCHOR.TOP)

    # Right top: evidence / why a KB article
    sh_why = _rect(slide, right_x, content_y, col_w, Inches(2.65), fill=WHITE, border=BORDER)
    _text_frame(sh_why, [
        ("Why this should become a KB article:", 12, True, NAVY, PP_ALIGN.LEFT),
        ("", 4, False, WHITE, PP_ALIGN.LEFT),
        (evidence[:280], 12, False, DARK, PP_ALIGN.LEFT),
    ], MSO_ANCHOR.TOP)

    # Right bottom: next step
    sh_next = _rect(slide, right_x, content_y + Inches(2.85), col_w, Inches(2.65),
                    fill=WHITE, border=GREEN)
    _text_frame(sh_next, [
        ("Next Step:", 12, True, GREEN, PP_ALIGN.LEFT),
        ("", 4, False, WHITE, PP_ALIGN.LEFT),
        (next_step[:320], 12, False, DARK, PP_ALIGN.LEFT),
    ], MSO_ANCHOR.TOP)

    # Speaker notes include MaturitySignal for presenter context
    maturity = last_occ.get("MaturitySignal") or "Unknown"
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = (
        f"Maturity Signal: {maturity}\n"
        f"Sessions seen: {occ_count} | First: {entry.get('first_seen', 'N/A')} | Last: {entry.get('last_seen', 'N/A')}\n"
        f"Evidence: {evidence[:400]}"
    )


def _slide_live_wins(prs: Presentation, delivered_rows: list[dict], current_date: str) -> None:
    """Live Wins This Week — only rendered when delivered_rows is non-empty."""
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, f"Live Wins This Week  \u2013  {current_date}")
    _rect(slide, Inches(0), Inches(1.4), Inches(0.18), H - Inches(1.4), fill=RGBColor(0x05, 0x96, 0x69))

    col_headers = ["Opportunity", "Process Stage", "Sub. Function", "Evidence"]
    col_x      = [Inches(0.35), Inches(5.0), Inches(7.2), Inches(9.4)]
    col_widths  = [Inches(4.55), Inches(2.1),  Inches(2.1),  Inches(3.75)]
    row_h = Inches(0.5)
    header_y = Inches(1.55)

    _rect(slide, Inches(0.35), header_y, W - Inches(0.7), row_h, fill=RGBColor(0x05, 0x96, 0x69))
    for hdr, cx, cw in zip(col_headers, col_x, col_widths):
        sh = _rect(slide, cx, header_y, cw, row_h)
        _text_frame(sh, [(hdr, 11, True, WHITE, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)

    for r_idx, row in enumerate(delivered_rows[:8]):
        y = header_y + row_h * (r_idx + 1)
        bg = WHITE if r_idx % 2 == 0 else ROW_ALT
        _rect(slide, Inches(0.35), y, W - Inches(0.7), row_h, fill=bg)
        row_vals = [
            (row.get("Title") or "")[:70],
            (row.get("ProcessStage") or "")[:28],
            (row.get("SubOrdinateFunction") or "")[:28],
            (row.get("EvidenceSummary") or "")[:80],
        ]
        for val, cx, cw in zip(row_vals, col_x, col_widths):
            sh = _rect(slide, cx, y, cw, row_h)
            _text_frame(sh, [(val, 10, False, DARK, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = (
        f"Live Wins — {len(delivered_rows)} opportunity/ies confirmed as Delivered / Active Today "
        f"in the {current_date} transcript.\n\n"
        + "\n".join(
            f"- {r.get('Title', '')}: {r.get('EvidenceSummary', '')[:120]}"
            for r in delivered_rows[:8]
        )
    )


def _slide_next_steps(prs: Presentation, near_term_asks: list[str]) -> None:
    """Last slide – Demo Intake / KB Strategy / Next Steps."""
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, "Demo Intake, KB Strategy, and Next Steps")

    mid = Inches(6.85)

    # Process flow (left column)
    sh_lbl_l = _rect(slide, Inches(0.4), Inches(1.55), Inches(6.0), Inches(0.42))
    _text_frame(sh_lbl_l, [("Process Flow", 13, True, NAVY, PP_ALIGN.LEFT)])

    flow_steps = [
        "Submit AI tools, prompts, agents, automations, reports, or workflows",
        "Identify SME / power users for each use case",
        "Present in a working session demo — rough work is welcome",
        "Classify into the AI Acceleration SharePoint framework",
        "Promote approved items to SharePoint (human review required)",
    ]
    sh_flow = _rect(slide, Inches(0.4), Inches(2.1), Inches(6.2), Inches(5.0),
                    fill=WHITE, border=BORDER)
    _bullet_frame(sh_flow, flow_steps, 13, DARK, prefix="\u2192  ")

    # Near-term asks (right column)
    sh_lbl_r = _rect(slide, mid, Inches(1.55), Inches(6.0), Inches(0.42))
    _text_frame(sh_lbl_r, [("Near-Term Ask", 13, True, NAVY, PP_ALIGN.LEFT)])

    asks = near_term_asks or [
        "Submit existing AI tools, prompts, agents, or workflows via the demo intake form",
        "Identify other power users in your team",
        "Share rough work \u2014 it does not need to be polished",
        "Review KB article candidates and confirm ownership",
    ]
    sh_asks = _rect(slide, mid, Inches(2.1), Inches(6.2), Inches(5.0),
                    fill=WHITE, border=BORDER)
    _bullet_frame(sh_asks, asks, 13, DARK)


# ── Data preparation helpers ──────────────────────────────────────────────────

def _session_number(history: list[dict]) -> int:
    """Current session number = count of unique meeting dates in history."""
    if not history:
        return 1
    dates = {o["date"] for e in history for o in e.get("occurrences", [])}
    return max(len(dates), 1)


def _kb_candidates(history: list[dict], top_n: int = 6) -> list[dict]:
    """Return top KB candidates: opportunities seen in 2+ sessions, sorted by recurrence."""
    candidates = [e for e in history if len(e.get("occurrences", [])) >= 2]
    candidates.sort(key=lambda e: len(e.get("occurrences", [])), reverse=True)
    return candidates[:top_n]


def _demo_backlog(last_rows: list[dict], history: list[dict], top_n: int = 8) -> list[dict]:
    """Build a demo backlog from last week's classified rows + history fill."""
    demos: list[dict] = []
    seen: set[str] = set()

    # Primary: last week's rows, sorted by ValueScore
    for row in sorted(last_rows or [], key=lambda r: -(r.get("ValueScore") or 0)):
        title = (row.get("Title") or "").strip()
        if title and title not in seen:
            demos.append({
                "owner": (row.get("SuggestedSMEChampionText")
                          or row.get("SuggestedBusinessOwnerText") or "TBD"),
                "title": title,
                "signal": row.get("SignalStrength", ""),
                "status": "Queued",
            })
            seen.add(title)
        if len(demos) >= top_n:
            break

    # Fill from history if needed
    for entry in sorted(history, key=lambda e: -len(e.get("occurrences", []))):
        if len(demos) >= top_n:
            break
        t = (entry.get("title") or "").strip()
        if t and t not in seen:
            last = entry["occurrences"][-1] if entry.get("occurrences") else {}
            demos.append({
                "owner": (last.get("SuggestedSMEChampionText")
                          or last.get("SuggestedBusinessOwnerText") or "TBD"),
                "title": t,
                "signal": last.get("SignalStrength", ""),
                "status": "Backlog",
            })
            seen.add(t)

    return demos[:top_n]


def _fallback_aligned_bullets(last_rows: list[dict]) -> list[str]:
    """Generate alignment bullets from last week's rows without LLM."""
    bullets: list[str] = []
    seen: set[str] = set()
    for row in sorted(last_rows or [], key=lambda r: -(r.get("ValueScore") or 0))[:6]:
        title = (row.get("Title") or "").strip()
        next_step = (row.get("NextStep") or "").strip()
        if title and title not in seen:
            b = f"{title} \u2014 {next_step[:80]}" if next_step else title
            bullets.append(b[:115])
            seen.add(title)
    return bullets or ["No prior session data available \u2014 update before meeting."]


def _default_agenda(session_num: int) -> list[str]:
    return [
        f"Recap last working session (Session {session_num - 1} highlights)",
        "AI working group / survey update",
        "Demo: [Presenter TBD \u2014 confirm before meeting]",
        "Reuse discussion: where else could this pattern apply?",
        "Review current KB article candidates",
        "Align on ownership, token costs, and next steps",
        "Decisions, owners, and next steps",
    ]


def _llm_slide_content(last_rows: list[dict], session_num: int) -> dict:
    """
    Call the LLM to generate aligned_on bullets, agenda items, and near-term asks.
    Falls back gracefully on any error.
    """
    try:
        from llm_client import call_llm
        skills_path = Path("skills") / "presentation_slide_content.md"
        system_prompt = (
            skills_path.read_text(encoding="utf-8") if skills_path.exists()
            else (
                "Generate PowerPoint slide content for the Electronics AI Working Group. "
                "Return only valid JSON: {\"aligned_on\": [], \"agenda_items\": [], \"near_term_asks\": []}"
            )
        )

        top_rows = sorted(last_rows or [], key=lambda r: -(r.get("ValueScore") or 0))[:8]
        summary = [
            {
                "title": r.get("Title", ""),
                "bucket": r.get("OperatingBucket", ""),
                "signal": r.get("SignalStrength", ""),
                "next_step": (r.get("NextStep") or "")[:120],
                "speaker": r.get("SourceSpeaker", ""),
                "evidence": (r.get("EvidenceSummary") or "")[:120],
            }
            for r in top_rows
        ]

        user_prompt = (
            f"Session number being recapped: {session_num}\n"
            f"Top opportunities from this session:\n"
            f"{json.dumps(summary, indent=2)}\n\n"
            "Generate slide content for the NEXT (upcoming) session. "
            'Return JSON: {"aligned_on": [4-5 bullets], '
            '"agenda_items": [6-8 items], "near_term_asks": [3-4 asks]}'
        )

        raw = call_llm(system_prompt, user_prompt)
        result = json.loads(raw)
        if isinstance(result.get("aligned_on"), list):
            return result
    except Exception as exc:
        print(f"[presentation] LLM content generation skipped ({type(exc).__name__}: {exc})")

    # Fallback: derive from data
    return {
        "aligned_on": _fallback_aligned_bullets(last_rows),
        "agenda_items": _default_agenda(session_num),
        "near_term_asks": [
            "Submit existing AI tools, prompts, agents, or workflows via the demo intake form",
            "Identify other power users in your team",
            "Review KB article candidates and confirm ownership",
            "Share rough work \u2014 it does not need to be polished",
        ],
    }


# ── Public API ────────────────────────────────────────────────────────────────

def build_meeting_presentation(
    history: list[dict],
    current_date: date,
    last_week_rows: list[dict] | None = None,
    out_dir: Path | None = None,
) -> Path:
    """
    Build the upcoming session PPTX from the current history and last week's rows.

    Args:
        history:         Full opportunity history list (after this week's ingest).
        current_date:    Meeting date of the transcript just processed.
        last_week_rows:  Classified rows from this week (used for slide content).
        out_dir:         Output directory; defaults to output/meeting_presentations/.

    Returns:
        Path to the generated .pptx file.
    """
    out_dir = Path(out_dir) if out_dir else PRES_DIR
    out_dir.mkdir(parents=True, exist_ok=True)

    session_num  = _session_number(history)       # the session just completed
    next_session = session_num + 1                # the session this deck prepares for
    next_date    = current_date + timedelta(days=7)

    print(f"[presentation] Building Session {next_session} deck for {next_date} ...")

    last_rows = last_week_rows or []
    content   = _llm_slide_content(last_rows, session_num)

    aligned_on     = content.get("aligned_on") or _fallback_aligned_bullets(last_rows)
    agenda_items   = content.get("agenda_items") or _default_agenda(next_session)
    near_term_asks = content.get("near_term_asks") or []

    kb_candidates = _kb_candidates(history)
    demos         = _demo_backlog(last_rows, history)

    # Build the deck
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    _slide_title(prs, next_session, next_date)
    _slide_confidential(prs, slide_num=2)
    _slide_agenda(prs, agenda_items, next_session)

    # Live Wins: only render if any delivered rows exist this week
    delivered = [r for r in last_rows if r.get("MaturitySignal") == "Delivered / Active Today"]
    if delivered:
        _slide_live_wins(prs, delivered, current_date.isoformat())

    _slide_aligned_on(prs, aligned_on, current_date.isoformat())
    _slide_demo_backlog(prs, demos)
    _slide_demo(prs, next_session)
    for i, kb in enumerate(kb_candidates, start=1):
        _slide_kb_candidate(prs, i, kb)
    _slide_next_steps(prs, near_term_asks)

    d        = next_date
    filename = f"Session_{next_session}_{d.month}_{d.day}_{d.year}.pptx"
    out_path = out_dir / filename

    # If the target file is open (e.g. in PowerPoint), write to a temp name.
    try:
        prs.save(str(out_path))
    except PermissionError:
        from datetime import datetime as _dt
        stamp = _dt.now().strftime("%H%M%S")
        alt_filename = f"Session_{next_session}_{d.month}_{d.day}_{d.year}_{stamp}.pptx"
        out_path = out_dir / alt_filename
        prs.save(str(out_path))
        print(f"[presentation] NOTE: original file was locked — saved as {alt_filename}")
        print(f"[presentation] Close the original in PowerPoint and rename if needed.")

    print(f"[presentation] saved -> {out_path}")
    return out_path
