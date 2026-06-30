"""
build_platform_vision_deck.py
Builds a simpler, leader-oriented overview deck that explains the vision:
turning the AI Transcript Intake Agent into a unified, config-driven
opportunity platform across business units (AI, Process Improvement, PMO, Sales).

Reuses the Acuity-branded shape helpers from src/presentation_builder.py and the
extra card/chip/divider helpers from demo/build_demo_deck.py so the deck matches
the look of the existing demo deck.

Output: demo/Unified_Opportunity_Platform_Vision.pptx

Run:
    .venv/Scripts/python.exe demo/build_platform_vision_deck.py
"""
from __future__ import annotations

import sys
from datetime import date
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "src"))
sys.path.insert(0, str(ROOT / "demo"))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import presentation_builder as pb
from presentation_builder import (
    NAVY, LIGHT_NAVY, TEAL, GOLD, GREEN, WHITE, DARK, MUTED, LIGHT,
    BORDER, ROW_ALT, W, H,
    _blank_slide, _slide_bg, _rect, _text_frame, _bullet_frame, _header_bar,
)

# Reuse the extra helpers already defined for the demo deck
from build_demo_deck import (
    RED, SUBTLE, _footer, _chip, _card, _section_divider,
)

OUT = Path(__file__).resolve().parent / "Unified_Opportunity_Platform_Vision.pptx"


# ── Slides ──────────────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = _blank_slide(prs)
    _slide_bg(slide, NAVY)
    _rect(slide, Inches(0), Inches(0), Inches(0.22), H, fill=TEAL)
    _rect(slide, Inches(0.22), Inches(0), W - Inches(0.22), Inches(1.05), fill=LIGHT_NAVY)

    sh_co = _rect(slide, Inches(0.5), Inches(0.1), W - Inches(1.0), Inches(0.85))
    _text_frame(sh_co, [(pb.COMPANY, 11, False, SUBTLE, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)

    sh_grp = _rect(slide, Inches(0.5), Inches(1.45), W - Inches(1.0), Inches(0.5))
    _text_frame(sh_grp, [("Enterprise AI & Opportunity Intake", 15, False, TEAL, PP_ALIGN.LEFT)],
                MSO_ANCHOR.MIDDLE)

    sh_title = _rect(slide, Inches(0.5), Inches(2.1), W - Inches(1.0), Inches(1.9))
    _text_frame(sh_title, [
        ("One Engine, Every Opportunity", 40, True, WHITE, PP_ALIGN.LEFT),
        ("Scaling the AI Intake Agent into a shared platform", 20, False, GOLD, PP_ALIGN.LEFT),
    ], MSO_ANCHOR.MIDDLE)

    sh_sub = _rect(slide, Inches(0.5), Inches(4.3), W - Inches(1.0), Inches(1.2))
    _bullet_frame(sh_sub, [
        "AI Opportunities  -  Process Improvement  -  PMO  -  Sales",
        "What it is  -  What unifies it  -  What we need from you  -  Roadmap",
    ], 14, SUBTLE, prefix="", gap_pt=6)

    sh_pres = _rect(slide, Inches(0.5), H - Inches(0.95), W - Inches(1.0), Inches(0.6))
    _text_frame(sh_pres, [
        ("Presented by David High (DK High)  -  " + date.today().strftime("%B %Y"),
         11, False, MUTED, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)


def slide_agenda(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, "What We'll Cover", "A simple walk from today's tool to the bigger vision")
    items = [
        ("1", "Where We Are Today", "A working AI opportunity intake pipeline"),
        ("2", "The Bigger Idea", "One engine, many business areas"),
        ("3", "How It Works", "Read, extract, classify, review, report"),
        ("4", "What Unifies It", "Shared backbone, your own language"),
        ("5", "The Data Picture", "Merge opportunities to find relationships"),
        ("6", "The Feedback Loop", "It personalizes to your team over time"),
        ("7", "What We Need From You", "Define your context, assign a reviewer"),
        ("8", "Roadmap", "From AI today to PMO and Sales later"),
    ]
    y = Inches(1.7)
    row_h = Inches(0.62)
    for n, title, desc in items:
        _chip(slide, Inches(0.7), y, Inches(0.55), Inches(0.5), n, TEAL, pt=16)
        sh = _rect(slide, Inches(1.45), y, W - Inches(2.1), Inches(0.5))
        _text_frame(sh, [(f"{title}   -   {desc}", 15, False, DARK, PP_ALIGN.LEFT)],
                    MSO_ANCHOR.MIDDLE)
        sh2 = _rect(slide, Inches(1.45), y, Inches(3.6), Inches(0.5))
        _text_frame(sh2, [(title, 15, True, NAVY, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)
        y = y + row_h
    _footer(slide, num)


def slide_today(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, "Where We Are Today", "A working pipeline, live now")

    sh = _rect(slide, Inches(0.7), Inches(1.55), W - Inches(1.4), Inches(0.9))
    _text_frame(sh, [(
        "We built a pipeline that turns meeting transcripts into structured, "
        "reviewable AI opportunities - with a human approving every step.",
        16, False, DARK, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)

    cards = [
        ("It works", [
            "Reads transcripts automatically",
            "Finds real opportunities with evidence",
            "Classifies them into a shared framework",
            "Live for the Electronics AI Working Group",
        ], GREEN),
        ("It's governed", [
            "Everything is draft-only",
            "Nothing is published without human review",
            "Evidence kept: speaker, timestamp, quote",
            "Safe by design",
        ], TEAL),
        ("It's tracked", [
            "Reports weekly and monthly",
            "Pushes approved items to Azure DevOps",
            "Syncs live status back into reports",
            "Answers 'where are we at with X?'",
        ], LIGHT_NAVY),
    ]
    cw = Inches(3.95)
    gap = Inches(0.2)
    x = Inches(0.7)
    for title, body, accent in cards:
        _card(slide, x, Inches(2.7), cw, Inches(3.5), title, body, accent=accent)
        x = x + cw + gap

    sh_q = _rect(slide, Inches(0.7), Inches(6.45), W - Inches(1.4), Inches(0.55), fill=NAVY)
    _text_frame(sh_q, [(
        "Today it serves one domain - AI opportunities. The same engine can serve many.",
        15, True, WHITE, PP_ALIGN.CENTER)], MSO_ANCHOR.MIDDLE)
    _footer(slide, num)


def slide_bigger_idea(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, "The Bigger Idea", "One core engine, tailored per business area")

    sh = _rect(slide, Inches(0.7), Inches(1.55), W - Inches(1.4), Inches(0.8))
    _text_frame(sh, [(
        "We don't rebuild for each team. We reuse one pipeline and plug in a "
        "'config pack' that fits each area's language and workflow.",
        15, False, DARK, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)

    # Core engine band
    sh_core = _rect(slide, Inches(0.7), Inches(2.6), W - Inches(1.4), Inches(0.8), fill=NAVY)
    _text_frame(sh_core, [(
        "CORE ENGINE (shared):  Read  >  Extract  >  Classify  >  Review  >  Report",
        15, True, WHITE, PP_ALIGN.CENTER)], MSO_ANCHOR.MIDDLE)

    # Domain packs row
    packs = [
        ("AI Opportunities", "Live today", GREEN),
        ("Process Improvement", "Next", GOLD),
        ("PMO", "Planned", LIGHT_NAVY),
        ("Sales", "Planned", LIGHT_NAVY),
    ]
    bw = Inches(2.85)
    gap = Inches(0.18)
    x = Inches(0.7)
    y = Inches(3.9)
    for t, d, c in packs:
        _chip(slide, x, y, bw, Inches(1.2), "", c)
        sh_t = _rect(slide, x, y + Inches(0.2), bw, Inches(0.5))
        _text_frame(sh_t, [(t, 15, True, WHITE, PP_ALIGN.CENTER)], MSO_ANCHOR.MIDDLE)
        sh_d = _rect(slide, x, y + Inches(0.68), bw, Inches(0.4))
        _text_frame(sh_d, [(d, 11, False, WHITE, PP_ALIGN.CENTER)], MSO_ANCHOR.MIDDLE)
        x = x + bw + gap

    sh_n = _rect(slide, Inches(0.7), Inches(5.7), W - Inches(1.4), Inches(0.8), fill=TEAL)
    _text_frame(sh_n, [(
        "Build once, reuse everywhere. Each area changes the vocabulary - not the engine.",
        15, True, WHITE, PP_ALIGN.CENTER)], MSO_ANCHOR.MIDDLE)
    _footer(slide, num)


def slide_how(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, "How It Works", "The same five steps for every business area")

    stages = [
        ("Read", "Source material in", LIGHT_NAVY),
        ("Extract", "Find the real opportunities", GOLD),
        ("Classify", "Map to your framework", GOLD),
        ("Review", "Human approves - never automatic", GREEN),
        ("Report", "Portfolio + progress", TEAL),
    ]
    bw = Inches(2.25)
    gap = Inches(0.18)
    x = Inches(0.7)
    y = Inches(2.4)
    for i, (t, d, c) in enumerate(stages):
        _chip(slide, x, y, bw, Inches(1.2), "", c)
        sh_t = _rect(slide, x, y + Inches(0.2), bw, Inches(0.5))
        _text_frame(sh_t, [(t, 16, True, WHITE, PP_ALIGN.CENTER)], MSO_ANCHOR.MIDDLE)
        sh_d = _rect(slide, x, y + Inches(0.66), bw, Inches(0.45))
        _text_frame(sh_d, [(d, 10, False, WHITE, PP_ALIGN.CENTER)], MSO_ANCHOR.MIDDLE)
        if i < len(stages) - 1:
            ar = _rect(slide, x + bw, y + Inches(0.47), gap, Inches(0.26))
            _text_frame(ar, [(">", 18, True, MUTED, PP_ALIGN.CENTER)], MSO_ANCHOR.MIDDLE)
        x = x + bw + gap

    _card(slide, Inches(0.7), Inches(4.1), Inches(5.9), Inches(2.3),
          "Deterministic where it should be", [
              "Reading, cleaning, validating, reporting is plain code",
              "Reproducible, fast, and free to run",
              "Same reliable backbone for every team",
          ], accent=TEAL)
    _card(slide, Inches(6.8), Inches(4.1), Inches(5.83), Inches(2.3),
          "AI only where it adds value", [
              "Understands language to find real opportunities",
              "Fills in the classification with nuance",
              "Always draft-only, always reviewed by a person",
          ], accent=GOLD)
    _footer(slide, num)


def slide_unifies(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, "What Unifies Every Area", "Your language on top of one shared backbone")

    _card(slide, Inches(0.7), Inches(1.7), Inches(5.9), Inches(4.6),
          "Shared backbone (the same everywhere)", [
              "Same evidence-based intake method",
              "Same scoring and governance guardrails",
              "Same human-review workflow",
              "Same portfolio reporting structure",
              "",
              "Result: leadership gets one comparable view across the enterprise.",
          ], accent=GREEN, body_pt=13)

    _card(slide, Inches(6.8), Inches(1.7), Inches(5.83), Inches(4.6),
          "Your configuration (unique per area)", [
              "Your process stages",
              "Your subordinate functions",
              "Your ownership roles",
              "Your priority signals",
              "",
              "Example: Product Vitality already exists as a workstream - it just "
              "needs its own process stages. Process Improvement needs its own config set.",
          ], accent=GOLD, body_pt=13)
    _footer(slide, num)


def slide_data_picture(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, "The Data Picture", "Separate intakes, one unified model")

    sh = _rect(slide, Inches(0.7), Inches(1.5), W - Inches(1.4), Inches(0.7))
    _text_frame(sh, [(
        "Each area keeps its own list. Underneath, shared dimensions let us merge "
        "them to find relationships, dependencies, and duplicates.",
        14, False, DARK, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)

    # Two source tables
    _card(slide, Inches(0.7), Inches(2.4), Inches(3.7), Inches(2.4),
          "AI Opportunities", [
              "One record per opportunity",
              "Evidence + classification",
              "Owner, team, status",
          ], accent=TEAL)
    _card(slide, Inches(4.6), Inches(2.4), Inches(3.7), Inches(2.4),
          "Process Improvements", [
              "One record per improvement",
              "Same shape, different domain",
              "Owner, team, status",
          ], accent=GOLD)

    # Merge target
    sh_m = _rect(slide, Inches(8.5), Inches(2.4), Inches(4.13), Inches(2.4), fill=NAVY)
    _text_frame(sh_m, [
        ("Unified Model", 16, True, WHITE, PP_ALIGN.CENTER),
        ("Shared dimensions:", 12, True, TEAL, PP_ALIGN.CENTER),
        ("Workstream - Stage - Function", 11, False, WHITE, PP_ALIGN.CENTER),
        ("Owner - Team - Tool", 11, False, WHITE, PP_ALIGN.CENTER),
    ], MSO_ANCHOR.MIDDLE)

    sh_n = _rect(slide, Inches(0.7), Inches(5.3), W - Inches(1.4), Inches(1.1), fill=TEAL)
    _text_frame(sh_n, [
        ("This is how we break down silos:", 15, True, WHITE, PP_ALIGN.CENTER),
        ("surface cross-functional relationships and recommend where work connects.",
         13, False, WHITE, PP_ALIGN.CENTER),
    ], MSO_ANCHOR.MIDDLE)
    _footer(slide, num)


def slide_feedback(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, "The Feedback Loop", "It gets more accurate the more you use it")

    stages = [
        ("Reviewer corrects", "Fixes class, stage, owner, next step", GOLD),
        ("Captured", "Corrections are recorded", LIGHT_NAVY),
        ("Promoted", "Rolled into your config & rules", TEAL),
        ("Personalized", "Output fits how your team works", GREEN),
    ]
    bw = Inches(2.85)
    gap = Inches(0.18)
    x = Inches(0.7)
    y = Inches(2.4)
    for i, (t, d, c) in enumerate(stages):
        _chip(slide, x, y, bw, Inches(1.3), "", c)
        sh_t = _rect(slide, x, y + Inches(0.2), bw, Inches(0.5))
        _text_frame(sh_t, [(t, 15, True, WHITE, PP_ALIGN.CENTER)], MSO_ANCHOR.MIDDLE)
        sh_d = _rect(slide, x, y + Inches(0.68), bw, Inches(0.5))
        _text_frame(sh_d, [(d, 10, False, WHITE, PP_ALIGN.CENTER)], MSO_ANCHOR.MIDDLE)
        if i < len(stages) - 1:
            ar = _rect(slide, x + bw, y + Inches(0.5), gap, Inches(0.26))
            _text_frame(ar, [(">", 18, True, MUTED, PP_ALIGN.CENTER)], MSO_ANCHOR.MIDDLE)
        x = x + bw + gap

    sh_n = _rect(slide, Inches(0.7), Inches(4.4), W - Inches(1.4), Inches(1.8), fill=WHITE,
                 border=BORDER)
    _rect(slide, Inches(0.7), Inches(4.4), W - Inches(1.4), Inches(0.12), fill=GREEN)
    sh_b = _rect(slide, Inches(1.0), Inches(4.75), W - Inches(2.0), Inches(1.3))
    _bullet_frame(sh_b, [
        "More accurate classifications over time",
        "Faster reviews as the system learns your patterns",
        "Output that increasingly reflects your team's real language and priorities",
    ], 14, DARK, gap_pt=4)
    _footer(slide, num)


def slide_need_from_you(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, "What We Need From You", "A small investment to make your area accurate")

    steps = [
        ("1", "A context owner", "Someone who knows your area to define stages, functions, roles, and priority signals."),
        ("2", "A reviewer", "Validates output and feeds corrections back into the loop."),
        ("3", "Source material", "Transcripts, notes, or intake lists we can pilot with."),
        ("4", "A pilot scope", "Pick one workflow to prove value before scaling."),
    ]
    y = Inches(1.9)
    rh = Inches(1.0)
    for n, t, d in steps:
        _chip(slide, Inches(0.8), y, Inches(0.6), Inches(0.6), n, TEAL, pt=18)
        st = _rect(slide, Inches(1.6), y, Inches(3.4), Inches(0.6))
        _text_frame(st, [(t, 16, True, NAVY, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)
        sd = _rect(slide, Inches(5.1), y, W - Inches(5.8), Inches(0.7))
        _text_frame(sd, [(d, 13, False, DARK, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)
        y = y + rh

    sh_n = _rect(slide, Inches(0.7), Inches(6.45), W - Inches(1.4), Inches(0.5), fill=GREEN)
    _text_frame(sh_n, [(
        "A few hours of your team's context up front - then the system does the heavy lifting.",
        13, True, WHITE, PP_ALIGN.CENTER)], MSO_ANCHOR.MIDDLE)
    _footer(slide, num)


def slide_roadmap(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, "Roadmap", "From AI today to an enterprise platform")

    rows = [
        ("Now", "AI Opportunities live", "Electronics AI Working Group running weekly", GREEN),
        ("Next", "Workstream-specific stages", "Start with Product Vitality stage model", GOLD),
        ("Then", "Process Improvement pack", "Add the BPI config set and pilot", GOLD),
        ("Later", "Relationship model", "Connect opportunities, recommend cross-functional work", TEAL),
        ("Future", "PMO, Sales, and beyond", "Onboard new areas through config, not code", LIGHT_NAVY),
    ]
    y = Inches(1.7)
    rh = Inches(0.92)
    for phase, title, desc, c in rows:
        _chip(slide, Inches(0.7), y, Inches(1.4), Inches(0.66), phase, c, pt=14)
        st = _rect(slide, Inches(2.3), y, Inches(4.2), Inches(0.66))
        _text_frame(st, [(title, 15, True, NAVY, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)
        sd = _rect(slide, Inches(6.6), y, W - Inches(7.3), Inches(0.66))
        _text_frame(sd, [(desc, 12, False, DARK, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)
        y = y + rh
    _footer(slide, num)


def slide_close(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, NAVY)
    _rect(slide, Inches(0), Inches(0), Inches(0.22), H, fill=TEAL)
    sh_t = _rect(slide, Inches(0.9), Inches(1.4), W - Inches(1.8), Inches(1.2))
    _text_frame(sh_t, [("Why It Matters & The Ask", 32, True, WHITE, PP_ALIGN.LEFT)],
                MSO_ANCHOR.MIDDLE)

    sh_b = _rect(slide, Inches(0.9), Inches(2.7), W - Inches(1.8), Inches(2.7))
    _bullet_frame(sh_b, [
        "Less manual tracking, more visibility into what we're working on",
        "Consistent governance across every business unit",
        "One source of truth for 'what are we doing and where are we at?'",
        "Cross-functional opportunities surface automatically",
        "We have a proven engine - now we scale it through configuration",
    ], 16, SUBTLE, prefix="-  ", gap_pt=8)

    sh_c = _rect(slide, Inches(0.9), H - Inches(1.5), W - Inches(1.8), Inches(0.8))
    _text_frame(sh_c, [
        ("The ask: pick a pilot area, assign a context owner + reviewer.   -   David High (DK High)",
         15, True, GOLD, PP_ALIGN.LEFT),
    ], MSO_ANCHOR.MIDDLE)
    _footer(slide, num)


def slide_thanks(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, NAVY)
    _rect(slide, Inches(0), Inches(0), Inches(0.22), H, fill=GOLD)
    sh = _rect(slide, Inches(0.9), Inches(2.8), W - Inches(1.8), Inches(1.6))
    _text_frame(sh, [
        ("Thank you", 44, True, WHITE, PP_ALIGN.LEFT),
        ("Questions & discussion", 20, False, TEAL, PP_ALIGN.LEFT),
    ], MSO_ANCHOR.MIDDLE)
    _footer(slide, num)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    n = [1]

    def nx():
        v = n[0]
        n[0] += 1
        return v

    slide_title(prs); nx()
    slide_agenda(prs, nx())
    _section_divider(prs, "PART 1", "Where we are today", nx())
    slide_today(prs, nx())
    slide_bigger_idea(prs, nx())
    _section_divider(prs, "PART 2", "How it works & what unifies it", nx())
    slide_how(prs, nx())
    slide_unifies(prs, nx())
    slide_data_picture(prs, nx())
    _section_divider(prs, "PART 3", "Why it gets better", nx())
    slide_feedback(prs, nx())
    _section_divider(prs, "PART 4", "What we need & the roadmap", nx())
    slide_need_from_you(prs, nx())
    slide_roadmap(prs, nx())
    slide_close(prs, nx())
    slide_thanks(prs, nx())

    prs.save(str(OUT))
    print(f"Saved {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
