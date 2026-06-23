"""
build_demo_deck.py
Builds the 20-30 minute demo + overview presentation for the
AI Transcript Intake Agent.

Reuses the Acuity-branded shape helpers from src/presentation_builder.py
so the deck matches the look of the auto-generated meeting decks.

Output: demo/AI_Transcript_Intake_Agent_Demo.pptx

Run:
    .venv/Scripts/python.exe demo/build_demo_deck.py
"""
from __future__ import annotations

import sys
from datetime import date
from pathlib import Path

# Make src/ importable so we can reuse the brand helpers
ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "src"))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import presentation_builder as pb
from presentation_builder import (
    NAVY, LIGHT_NAVY, TEAL, GOLD, GREEN, WHITE, DARK, MUTED, LIGHT,
    BORDER, ROW_ALT, W, H,
    _blank_slide, _slide_bg, _rect, _text_frame, _bullet_frame, _header_bar,
)

RED = RGBColor(0xC0, 0x60, 0x4D)
SUBTLE = RGBColor(0xB9, 0xC4, 0xD1)

OUT = Path(__file__).resolve().parent / "AI_Transcript_Intake_Agent_Demo.pptx"


# ── Extra helpers ──────────────────────────────────────────────────────────────

def _footer(slide, num: int) -> None:
    sh = _rect(slide, Inches(0.5), H - Inches(0.45), W - Inches(1.0), Inches(0.35))
    _text_frame(sh, [(
        f"Electronics AI Working Group  |  AI Transcript Intake Agent  |  "
        f"Proprietary & Confidential  |  NYSE: AYI  |  {num}",
        9, False, MUTED, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)


def _chip(slide, x, y, w, h, label, fill, text_color=WHITE, pt=11, bold=True):
    sh = _rect(slide, x, y, w, h, fill=fill)
    _text_frame(sh, [(label, pt, bold, text_color, PP_ALIGN.CENTER)], MSO_ANCHOR.MIDDLE)
    return sh


def _card(slide, x, y, w, h, title, body_lines, accent=TEAL,
          title_pt=14, body_pt=11):
    """A white card with a colored top accent bar, title, and bullet body."""
    _rect(slide, x, y, w, h, fill=WHITE, border=BORDER)
    _rect(slide, x, y, w, Inches(0.12), fill=accent)
    th = Inches(0.5)
    sh_t = _rect(slide, x + Inches(0.18), y + Inches(0.2), w - Inches(0.36), th)
    _text_frame(sh_t, [(title, title_pt, True, NAVY, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)
    sh_b = _rect(slide, x + Inches(0.18), y + Inches(0.78),
                 w - Inches(0.36), h - Inches(0.95))
    _bullet_frame(sh_b, body_lines, body_pt, DARK, gap_pt=3)


def _section_divider(prs, kicker, title, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, NAVY)
    _rect(slide, Inches(0), Inches(0), Inches(0.22), H, fill=TEAL)
    sh_k = _rect(slide, Inches(0.9), Inches(2.7), W - Inches(1.8), Inches(0.6))
    _text_frame(sh_k, [(kicker, 16, False, TEAL, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)
    sh_t = _rect(slide, Inches(0.9), Inches(3.2), W - Inches(1.8), Inches(1.4))
    _text_frame(sh_t, [(title, 36, True, WHITE, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)
    _footer(slide, num)
    return slide


# ── Slides ──────────────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = _blank_slide(prs)
    _slide_bg(slide, NAVY)
    _rect(slide, Inches(0), Inches(0), Inches(0.22), H, fill=TEAL)
    _rect(slide, Inches(0.22), Inches(0), W - Inches(0.22), Inches(1.05), fill=LIGHT_NAVY)

    sh_co = _rect(slide, Inches(0.5), Inches(0.1), W - Inches(1.0), Inches(0.85))
    _text_frame(sh_co, [(pb.COMPANY, 11, False, SUBTLE, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)

    sh_grp = _rect(slide, Inches(0.5), Inches(1.45), W - Inches(1.0), Inches(0.5))
    _text_frame(sh_grp, [("Electronics AI Working Group", 15, False, TEAL, PP_ALIGN.LEFT)],
                MSO_ANCHOR.MIDDLE)

    sh_title = _rect(slide, Inches(0.5), Inches(2.1), W - Inches(1.0), Inches(1.9))
    _text_frame(sh_title, [
        ("AI Transcript Intake Agent", 40, True, WHITE, PP_ALIGN.LEFT),
        ("From meeting conversation to tracked execution", 20, False, GOLD, PP_ALIGN.LEFT),
    ], MSO_ANCHOR.MIDDLE)

    sh_sub = _rect(slide, Inches(0.5), Inches(4.3), W - Inches(1.0), Inches(1.2))
    _bullet_frame(sh_sub, [
        "Why we built it  -  How it works  -  The benefits  -  How it scales",
        "Live demo + overview  -  20-30 minutes",
    ], 14, SUBTLE, prefix="", gap_pt=6)

    sh_pres = _rect(slide, Inches(0.5), H - Inches(0.95), W - Inches(1.0), Inches(0.6))
    _text_frame(sh_pres, [
        ("Presented by David High (DK High)  -  " + date.today().strftime("%B %Y"),
         11, False, MUTED, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)


def slide_agenda(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, "Agenda", "What we'll cover in the next 20-30 minutes")
    items = [
        ("1", "The Problem", "Why this was created - the manual pain it removes"),
        ("2", "The Solution", "What the agent does, end to end"),
        ("3", "How It Works", "The pipeline, the AI, and the guardrails"),
        ("4", "The Outputs", "Weekly + monthly reports, ADO work items"),
        ("5", "Closing the Loop", "From discussion to tracked execution in Azure DevOps"),
        ("6", "Live Demo", "Drop a transcript, run it, see the report"),
        ("7", "The Benefits", "Time savings and ROI"),
        ("8", "How It Scales", "Other teams, other extraction types"),
    ]
    y = Inches(1.7)
    row_h = Inches(0.62)
    for n, title, desc in items:
        _chip(slide, Inches(0.7), y, Inches(0.55), Inches(0.5), n, TEAL, pt=16)
        sh = _rect(slide, Inches(1.45), y, W - Inches(2.1), Inches(0.5))
        _text_frame(sh, [
            (f"{title}   -   {desc}", 15, False, DARK, PP_ALIGN.LEFT),
        ], MSO_ANCHOR.MIDDLE)
        # bold the title portion by overlaying
        sh2 = _rect(slide, Inches(1.45), y, Inches(3.2), Inches(0.5))
        _text_frame(sh2, [(title, 15, True, NAVY, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)
        y = y + row_h
    _footer(slide, num)


def slide_problem(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, "The Problem", "Why we built this")

    sh = _rect(slide, Inches(0.7), Inches(1.55), W - Inches(1.4), Inches(0.9))
    _text_frame(sh, [
        ("Every week the working group meets and talks about ways the business "
         "could use AI. Those ideas were getting lost.", 16, False, DARK, PP_ALIGN.LEFT),
    ], MSO_ANCHOR.MIDDLE)

    cards = [
        ("Manual & slow", [
            "Someone had to read hours of transcript",
            "Pull out the real opportunities by hand",
            "Fill ~40 fields per idea in a spreadsheet",
            "~7 hours of effort, every single week",
        ], RED),
        ("Inconsistent", [
            "Different people captured different things",
            "No shared classification language",
            "Hard to compare week over week",
            "Quality depended on who did it",
        ], GOLD),
        ("No follow-through", [
            "Ideas were captured, then forgotten",
            "No link to who would actually do them",
            "'Did anything happen with X?' - nobody knew",
            "The loop was never closed",
        ], LIGHT_NAVY),
    ]
    cw = Inches(3.95)
    gap = Inches(0.2)
    x = Inches(0.7)
    for title, body, accent in cards:
        _card(slide, x, Inches(2.7), cw, Inches(3.5), title, body, accent=accent)
        x = x + cw + gap

    sh_q = _rect(slide, Inches(0.7), Inches(6.45), W - Inches(1.4), Inches(0.55), fill=NAVY)
    _text_frame(sh_q, [(
        "The core insight: a transcript is the BEGINNING of the work, not the end of it.",
        15, True, WHITE, PP_ALIGN.CENTER)], MSO_ANCHOR.MIDDLE)
    _footer(slide, num)


def slide_solution(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, "The Solution", "A local agent that turns talk into tracked work")

    sh = _rect(slide, Inches(0.7), Inches(1.55), W - Inches(1.4), Inches(0.8))
    _text_frame(sh, [(
        "Drop in a meeting transcript. Get back reviewable AI-opportunity records, "
        "trend reports, and tracked work items - with a human in the loop at every step.",
        15, False, DARK, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)

    # horizontal flow of 5 stages
    stages = [
        ("Read", "DOCX transcript in", LIGHT_NAVY),
        ("Extract", "AI finds opportunities", GOLD),
        ("Classify", "Map to framework", GOLD),
        ("Report", "Weekly + monthly HTML", TEAL),
        ("Track", "Push to Azure DevOps", GREEN),
    ]
    bw = Inches(2.25)
    gap = Inches(0.18)
    x = Inches(0.7)
    y = Inches(2.6)
    for i, (t, d, c) in enumerate(stages):
        _chip(slide, x, y, bw, Inches(1.1), "", c)
        sh_t = _rect(slide, x, y + Inches(0.15), bw, Inches(0.5))
        _text_frame(sh_t, [(t, 16, True, WHITE, PP_ALIGN.CENTER)], MSO_ANCHOR.MIDDLE)
        sh_d = _rect(slide, x, y + Inches(0.6), bw, Inches(0.4))
        _text_frame(sh_d, [(d, 10, False, WHITE, PP_ALIGN.CENTER)], MSO_ANCHOR.MIDDLE)
        if i < len(stages) - 1:
            ar = _rect(slide, x + bw, y + Inches(0.42), gap, Inches(0.26))
            _text_frame(ar, [(">", 18, True, MUTED, PP_ALIGN.CENTER)], MSO_ANCHOR.MIDDLE)
        x = x + bw + gap

    # principles row
    _card(slide, Inches(0.7), Inches(4.2), Inches(5.9), Inches(2.3),
          "What makes it trustworthy", [
              "Full transcripts are NEVER sent to the model - only filtered chunks",
              "Every AI output is DRAFT-only, defaults to 'Needs Review'",
              "Human review required before any SharePoint or ADO push",
              "Evidence preserved: speaker, timestamp, quote on every row",
          ], accent=GREEN)
    _card(slide, Inches(6.8), Inches(4.2), Inches(5.83), Inches(2.3),
          "What makes it efficient", [
              "Deterministic Python does all the heavy lifting (free)",
              "AI used only for the two semantic steps: extract + classify",
              "Runs on your existing GitHub Copilot subscription - no new key",
              "Idempotent: safe to re-run, never double-counts",
          ], accent=TEAL)
    _footer(slide, num)


def slide_how_pipeline(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, "How It Works - The Pipeline", "Mostly deterministic Python; AI only where it adds value")

    rows = [
        ("0", "ADO status sync", "Pull live work-item status into reports", "Python", LIGHT_NAVY),
        ("1-3", "Read - Clean - Chunk", "Parse DOCX, keep speaker turns, keyword-filter", "Python", LIGHT_NAVY),
        ("4", "Extract candidates", "Find AI opportunities in filtered chunks", "AI", GOLD),
        ("5", "Classify", "Map each to the 40-field framework", "AI", GOLD),
        ("5b", "Validate", "Check choices, confidence-filter", "Python", LIGHT_NAVY),
        ("6-8", "Export", "review_rows.xlsx, payload, HTML report", "Python", TEAL),
        ("9-11", "Archive + Report", "History store, weekly + monthly reports", "Python", TEAL),
        ("9c", "ADO push (opt-in)", "Create tracked Issues after review", "Python", GREEN),
    ]
    y = Inches(1.55)
    rh = Inches(0.58)
    # header
    hdr = _rect(slide, Inches(0.7), y, W - Inches(1.4), Inches(0.45), fill=NAVY)
    _text_frame(hdr, [("Step        Stage                         What happens                                                     Type",
                       11, True, WHITE, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)
    y = y + Inches(0.5)
    for i, (step, stage, what, typ, c) in enumerate(rows):
        bg = ROW_ALT if i % 2 else WHITE
        _rect(slide, Inches(0.7), y, W - Inches(1.4), rh, fill=bg, border=BORDER)
        s1 = _rect(slide, Inches(0.85), y, Inches(0.9), rh)
        _text_frame(s1, [(step, 12, True, NAVY, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)
        s2 = _rect(slide, Inches(1.8), y, Inches(3.0), rh)
        _text_frame(s2, [(stage, 12, True, DARK, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)
        s3 = _rect(slide, Inches(4.9), y, Inches(6.0), rh)
        _text_frame(s3, [(what, 11, False, DARK, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)
        _chip(slide, Inches(11.3), y + Inches(0.1), Inches(1.1), Inches(0.42),
              typ, c, pt=10)
        y = y + rh
    _footer(slide, num)


def slide_ai_vs_det(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, "How It Works - AI vs. Deterministic", "Two AI steps. Everything else is plain code.")

    _card(slide, Inches(0.7), Inches(1.7), Inches(5.9), Inches(4.6),
          "AI does the SEMANTIC work (2 steps)", [
              "Extraction: reads filtered chunks and asks 'is this a real AI opportunity?'",
              "Classification: fills ~40 framework fields per opportunity",
              "One model call per chunk / per candidate - small, cheap prompts",
              "Runs on GitHub Copilot SDK by default (or OpenAI if a key is set)",
              "",
              "Why AI here? Judgement, nuance, and language understanding that rules can't capture.",
          ], accent=GOLD, body_pt=12)

    _card(slide, Inches(6.8), Inches(1.7), Inches(5.83), Inches(4.6),
          "Python does EVERYTHING else (deterministic)", [
              "DOCX reading, cleaning, speaker-turn splitting",
              "Keyword filtering and chunking (token control)",
              "JSON validation against allowed choice values",
              "Excel export, payload building, HTML report generation",
              "History store, trend analysis, dedup matching",
              "ADO REST API calls (push + status sync)",
              "",
              "Why Python here? Reproducible, free, fast, and testable.",
          ], accent=TEAL, body_pt=12)
    _footer(slide, num)


def slide_framework(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, "The Classification Framework", "A shared language for every opportunity")

    sh = _rect(slide, Inches(0.7), Inches(1.5), W - Inches(1.4), Inches(0.6))
    _text_frame(sh, [(
        "Each opportunity is mapped onto consistent dimensions so they can be "
        "compared, prioritized, and tracked over time.", 14, False, DARK, PP_ALIGN.LEFT)],
        MSO_ANCHOR.MIDDLE)

    dims = [
        ("AI Use Case Type", "Classification - Action - Both", TEAL),
        ("Operating Bucket", "Pre-Sale, Manufacturing, Post-Shipment, Governance...", LIGHT_NAVY),
        ("Level of Analysis", "0 Signal -> 7 Release Candidate", GOLD),
        ("Maturity Signal", "Aspirational -> Exploring -> Piloting -> Deployed", GREEN),
        ("Signal Strength", "Isolated -> Cross-Functional -> Leadership Priority", TEAL),
        ("Value / Effort / Risk", "1-5 scores for prioritization", LIGHT_NAVY),
    ]
    cw = Inches(3.95)
    ch = Inches(1.55)
    gap = Inches(0.2)
    x = Inches(0.7)
    y = Inches(2.4)
    for i, (t, d, c) in enumerate(dims):
        col = i % 3
        row = i // 3
        cx = Inches(0.7) + col * (cw + gap)
        cy = y + row * (ch + Inches(0.25))
        _rect(slide, cx, cy, cw, ch, fill=WHITE, border=BORDER)
        _rect(slide, cx, cy, Inches(0.12), ch, fill=c)
        st = _rect(slide, cx + Inches(0.3), cy + Inches(0.18), cw - Inches(0.4), Inches(0.5))
        _text_frame(st, [(t, 14, True, NAVY, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)
        sd = _rect(slide, cx + Inches(0.3), cy + Inches(0.7), cw - Inches(0.4), Inches(0.7))
        _text_frame(sd, [(d, 11, False, DARK, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)

    sh_n = _rect(slide, Inches(0.7), Inches(6.4), W - Inches(1.4), Inches(0.55), fill=NAVY)
    _text_frame(sh_n, [(
        "The framework is just config (JSON). Swap it, and the same engine classifies a "
        "completely different domain.", 14, True, WHITE, PP_ALIGN.CENTER)], MSO_ANCHOR.MIDDLE)
    _footer(slide, num)


def slide_outputs(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, "The Outputs", "Self-contained, shareable, no server needed")

    _card(slide, Inches(0.7), Inches(1.6), Inches(5.9), Inches(2.45),
          "Weekly report (5 tabs)", [
              "Cards - one per opportunity, with ADO status chip",
              "Analytics - charts by bucket, type, level, maturity",
              "Full Table - all 40 fields, sortable",
              "Trends - new vs. carried-over, week-over-week movement",
              "Progress - live ADO status, grouped by state",
          ], accent=TEAL)
    _card(slide, Inches(6.8), Inches(1.6), Inches(5.83), Inches(2.45),
          "Monthly rollup", [
              "Unique opportunity count, weeks covered",
              "New vs. carried-from-prior month",
              "Per-week bucket distribution",
              "Signal / level momentum (escalations)",
              "Full table with first-seen / last-seen",
          ], accent=GOLD)
    _card(slide, Inches(0.7), Inches(4.25), Inches(5.9), Inches(2.1),
          "For review & hand-off", [
              "review_rows.xlsx - human review workbook",
              "sharepoint_payload.json - draft SharePoint rows",
              "master_opportunities.xlsx - Power BI source (incl. ADO columns)",
          ], accent=GREEN)
    _card(slide, Inches(6.8), Inches(4.25), Inches(5.83), Inches(2.1),
          "For the meeting", [
              "Auto-generated PPTX deck for the upcoming session",
              "Demo backlog, KB candidates, next steps",
              "Built from accumulated history, not by hand",
          ], accent=LIGHT_NAVY)
    _footer(slide, num)


def slide_ado(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, "Closing the Loop - Azure DevOps", "From 'we discussed it' to 'someone owns it'")

    sh = _rect(slide, Inches(0.7), Inches(1.5), W - Inches(1.4), Inches(0.75))
    _text_frame(sh, [(
        "The newest capability. Identified opportunities only deliver value if someone "
        "picks them up. ADO integration was added to close that gap.", 15, False, DARK,
        PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)

    _card(slide, Inches(0.7), Inches(2.45), Inches(3.9), Inches(3.0),
          "Push (opt-in)", [
              "After review, run with --push-ado",
              "Each approved opportunity becomes a tracked ADO Issue",
              "Rich description: problem, evidence, next step, metadata",
              "Title-based dedup - safe to re-run",
          ], accent=GREEN)
    _card(slide, Inches(4.75), Inches(2.45), Inches(3.9), Inches(3.0),
          "Sync (automatic)", [
              "Every weekly run pulls live ADO status first",
              "To Do / Doing / Done, assignee, iteration",
              "Read-only - never creates or edits items",
              "Silent if ADO isn't configured",
          ], accent=TEAL)
    _card(slide, Inches(8.8), Inches(2.45), Inches(3.83), Inches(3.0),
          "See it (Progress tab)", [
              "Every opportunity grouped by current state",
              "Highlights items that moved in last 7 days",
              "Direct links to the ADO work items",
              "Answers 'where are we at with X?'",
          ], accent=GOLD)

    sh_q = _rect(slide, Inches(0.7), Inches(5.7), W - Inches(1.4), Inches(0.75), fill=NAVY)
    _text_frame(sh_q, [(
        "This is the difference between a meeting note and a managed pipeline of work.",
        16, True, WHITE, PP_ALIGN.CENTER)], MSO_ANCHOR.MIDDLE)
    _footer(slide, num)


def slide_demo(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, NAVY)
    _rect(slide, Inches(0), Inches(0), Inches(0.22), H, fill=GREEN)
    sh_k = _rect(slide, Inches(0.9), Inches(1.7), W - Inches(1.8), Inches(0.6))
    _text_frame(sh_k, [("LIVE DEMO", 18, True, GREEN, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)
    sh_t = _rect(slide, Inches(0.9), Inches(2.25), W - Inches(1.8), Inches(1.0))
    _text_frame(sh_t, [("Let's run it.", 38, True, WHITE, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)

    sh_b = _rect(slide, Inches(0.9), Inches(3.5), W - Inches(1.8), Inches(3.0))
    _bullet_frame(sh_b, [
        "1.  Drop a DOCX transcript into input/transcripts/",
        "2.  Run:  python src/main.py --input \"...\" --mode weekly",
        "3.  Open the weekly report - walk the 5 tabs",
        "4.  Review a few cards (everything is 'Needs Review')",
        "5.  Re-run with --push-ado  ->  watch ADO Issues get created",
        "6.  Open the Progress tab - live status, linked to ADO",
    ], 16, SUBTLE, prefix="", gap_pt=8)
    _footer(slide, num)


def slide_benefits(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, "The Benefits - Time & ROI", "Measured against doing it all by hand")

    kpis = [
        ("7.5 h", "Saved per week", GREEN),
        ("83%", "Less manual effort", TEAL),
        ("390 h", "Saved per year", GOLD),
        ("~$33k", "Annual time value*", LIGHT_NAVY),
    ]
    cw = Inches(2.95)
    gap = Inches(0.18)
    x = Inches(0.7)
    for val, lbl, c in kpis:
        _rect(slide, x, Inches(1.6), cw, Inches(1.6), fill=WHITE, border=BORDER)
        _rect(slide, x, Inches(1.6), cw, Inches(0.12), fill=c)
        sv = _rect(slide, x, Inches(1.85), cw, Inches(0.8))
        _text_frame(sv, [(val, 32, True, NAVY, PP_ALIGN.CENTER)], MSO_ANCHOR.MIDDLE)
        sl = _rect(slide, x, Inches(2.6), cw, Inches(0.5))
        _text_frame(sl, [(lbl, 12, False, DARK, PP_ALIGN.CENTER)], MSO_ANCHOR.MIDDLE)
        x = x + cw + gap

    _card(slide, Inches(0.7), Inches(3.5), Inches(5.9), Inches(2.9),
          "Transcript processing (~5.8 h/wk)", [
              "Read + parse transcript: 1.8 h -> 0.1 h",
              "Identify opportunities: 1.5 h -> 0.2 h",
              "Classify into 40-field schema: 2.0 h -> 0.2 h",
              "Build workbook + payload + reports: 1.0 h -> 0.1 h",
              "QA and review: largely retained (human judgement)",
          ], accent=TEAL)
    _card(slide, Inches(6.8), Inches(3.5), Inches(5.83), Inches(2.9),
          "ADO lifecycle (~1.7 h/wk)", [
              "Work item creation: ~2.3 h/cycle -> minutes",
              "Status tracking + follow-up: ~1 h/wk -> ~5 min",
              "Traceability lookups: ~20 min/wk -> 0",
              "Dedup-safe re-runs eliminate rework risk",
              "No AI credits - pure REST API",
          ], accent=GREEN)

    sh_n = _rect(slide, Inches(0.7), Inches(6.55), W - Inches(1.4), Inches(0.4))
    _text_frame(sh_n, [(
        "* Optional labor-value lens at $85/hr loaded rate. Recurring runtime: ~30 Copilot "
        "credits per transcript.", 10, False, MUTED, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)
    _footer(slide, num)


def slide_scale(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, "How It Scales", "Same engine, different domains")

    sh = _rect(slide, Inches(0.7), Inches(1.5), W - Inches(1.4), Inches(0.7))
    _text_frame(sh, [(
        "The engine is domain-agnostic. The framework lives in config (JSON) and the "
        "extraction logic in skills (Markdown). Swap those, point it at different source "
        "documents, and it works for a new team.", 14, False, DARK, PP_ALIGN.LEFT)],
        MSO_ANCHOR.MIDDLE)

    examples = [
        ("Sales / CX", "Call transcripts", "Extract customer pain points, feature requests, churn signals", TEAL),
        ("HR / People", "Interview & survey notes", "Extract themes, sentiment, policy gaps, training needs", GOLD),
        ("Engineering", "Incident retros", "Extract action items, root causes, recurring failure modes", RED),
        ("Product", "User research sessions", "Extract jobs-to-be-done, usability issues, prioritized asks", LIGHT_NAVY),
        ("Legal / Risk", "Contract reviews", "Extract obligations, risks, non-standard clauses", GREEN),
        ("Operations", "Standups & reviews", "Extract blockers, dependencies, process improvements", TEAL),
    ]
    cw = Inches(3.95)
    ch = Inches(1.5)
    gap = Inches(0.2)
    y0 = Inches(2.4)
    for i, (team, src, what, c) in enumerate(examples):
        col = i % 3
        row = i // 3
        cx = Inches(0.7) + col * (cw + gap)
        cy = y0 + row * (ch + Inches(0.22))
        _rect(slide, cx, cy, cw, ch, fill=WHITE, border=BORDER)
        _rect(slide, cx, cy, cw, Inches(0.1), fill=c)
        st = _rect(slide, cx + Inches(0.22), cy + Inches(0.18), cw - Inches(0.4), Inches(0.4))
        _text_frame(st, [(f"{team}", 14, True, NAVY, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)
        ss = _rect(slide, cx + Inches(0.22), cy + Inches(0.58), cw - Inches(0.4), Inches(0.35))
        _text_frame(ss, [(src, 11, True, c, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)
        sw = _rect(slide, cx + Inches(0.22), cy + Inches(0.9), cw - Inches(0.4), Inches(0.55))
        _text_frame(sw, [(what, 10, False, DARK, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)

    sh_n = _rect(slide, Inches(0.7), Inches(6.45), W - Inches(1.4), Inches(0.5), fill=NAVY)
    _text_frame(sh_n, [(
        "What stays the same: ingestion, chunking, draft-only guardrails, human review, "
        "trend reporting, and the ADO loop.", 13, True, WHITE, PP_ALIGN.CENTER)],
        MSO_ANCHOR.MIDDLE)
    _footer(slide, num)


def slide_adopt(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, LIGHT)
    _header_bar(slide, "Adopting It For Your Team", "What it takes to stand up a new instance")

    steps = [
        ("1", "Define the framework", "List the dimensions and choice values for your domain in config JSON."),
        ("2", "Write the extraction prompt", "One skill file describing what an 'opportunity' looks like for you."),
        ("3", "Point at your sources", "Drop your transcripts / notes / docs into the input folder."),
        ("4", "Run & review", "The same pipeline produces drafts; a human reviews them."),
        ("5", "Wire your tracker", "Map outputs to ADO, Jira, SharePoint - whatever closes your loop."),
    ]
    y = Inches(1.8)
    rh = Inches(0.92)
    for n, t, d in steps:
        _chip(slide, Inches(0.8), y, Inches(0.6), Inches(0.6), n, TEAL, pt=18)
        st = _rect(slide, Inches(1.6), y, Inches(3.4), Inches(0.6))
        _text_frame(st, [(t, 16, True, NAVY, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)
        sd = _rect(slide, Inches(5.1), y, W - Inches(5.8), Inches(0.6))
        _text_frame(sd, [(d, 13, False, DARK, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)
        y = y + rh

    sh_n = _rect(slide, Inches(0.7), Inches(6.45), W - Inches(1.4), Inches(0.5), fill=GREEN)
    _text_frame(sh_n, [(
        "No new infrastructure. Runs locally on an existing Copilot subscription. "
        "Governed by human review end to end.", 13, True, WHITE, PP_ALIGN.CENTER)],
        MSO_ANCHOR.MIDDLE)
    _footer(slide, num)


def slide_close(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, NAVY)
    _rect(slide, Inches(0), Inches(0), Inches(0.22), H, fill=TEAL)
    sh_t = _rect(slide, Inches(0.9), Inches(1.6), W - Inches(1.8), Inches(1.2))
    _text_frame(sh_t, [("Recap & Next Steps", 34, True, WHITE, PP_ALIGN.LEFT)], MSO_ANCHOR.MIDDLE)

    sh_b = _rect(slide, Inches(0.9), Inches(2.9), W - Inches(1.8), Inches(2.7))
    _bullet_frame(sh_b, [
        "Turns weekly meeting talk into structured, classified, reviewable opportunities",
        "Saves ~7.5 hours/week (~390 hours/year) versus doing it by hand",
        "Closes the loop with Azure DevOps - discussion becomes tracked execution",
        "Human-in-the-loop and draft-only at every step - safe by design",
        "Domain-agnostic - swap the config to serve any team or extraction type",
    ], 16, SUBTLE, prefix="-  ", gap_pt=8)

    sh_c = _rect(slide, Inches(0.9), H - Inches(1.5), W - Inches(1.8), Inches(0.8))
    _text_frame(sh_c, [
        ("Want it for your team? Let's talk.   -   David High (DK High)",
         16, True, GOLD, PP_ALIGN.LEFT),
    ], MSO_ANCHOR.MIDDLE)
    _footer(slide, num)


def slide_thanks(prs, num):
    slide = _blank_slide(prs)
    _slide_bg(slide, NAVY)
    _rect(slide, Inches(0), Inches(0), Inches(0.22), H, fill=GOLD)
    sh = _rect(slide, Inches(0.9), Inches(2.8), W - Inches(1.8), Inches(1.6))
    _text_frame(sh, [
        ("Thank you", 44, True, WHITE, PP_ALIGN.LEFT),
        ("Questions & discussion", 20, False, TEAL, PP_ALIGN.LEFT),
    ], MSO_ANCHOR.MIDDLE)
    _footer(slide, num)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    n = [1]

    def nx():
        v = n[0]
        n[0] += 1
        return v

    slide_title(prs); nx()
    slide_agenda(prs, nx())
    _section_divider(prs, "PART 1", "Why we built it", nx())
    slide_problem(prs, nx())
    slide_solution(prs, nx())
    _section_divider(prs, "PART 2", "How it works", nx())
    slide_how_pipeline(prs, nx())
    slide_ai_vs_det(prs, nx())
    slide_framework(prs, nx())
    slide_outputs(prs, nx())
    slide_ado(prs, nx())
    _section_divider(prs, "PART 3", "See it in action", nx())
    slide_demo(prs, nx())
    _section_divider(prs, "PART 4", "Benefits & scale", nx())
    slide_benefits(prs, nx())
    slide_scale(prs, nx())
    slide_adopt(prs, nx())
    slide_close(prs, nx())
    slide_thanks(prs, nx())

    prs.save(str(OUT))
    print(f"Saved {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
